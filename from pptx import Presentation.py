from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import math

prs = Presentation()

slides = [
    ("Deepfake & AI-Powered Cybercrime", "Understanding Modern Digital Threats"),
    ("What is Deepfake?", "AI-generated fake videos, voices, and images that look real."),
    ("AI Cybercrime", "Hackers use AI for scams, phishing, malware, and identity theft."),
    ("Voice Cloning Scam", "Criminals clone voices to trick people into sending money."),
    ("Detection & Safety", "Verify media, use 2FA, and stay aware of online scams."),
    ("Conclusion", "Cybersecurity awareness is essential in the AI era.")
]

# Create modern cyber-style images
image_paths = []

for i, (title, _) in enumerate(slides, start=1):
    img = Image.new("RGB", (1280, 720), (8, 12, 25))
    draw = ImageDraw.Draw(img)

    # Futuristic circles
    for r in range(50, 350, 40):
        draw.ellipse((640-r, 360-r, 640+r, 360+r), outline=(0, 255, 200))

    # Grid lines
    for x in range(0, 1280, 60):
        draw.line((x, 0, x, 720), fill=(20, 40, 60))
    for y in range(0, 720, 60):
        draw.line((0, y, 1280, y), fill=(20, 40, 60))

    # Main text
    draw.text((80, 120), title, fill=(255, 255, 255))
    draw.text((80, 220), "AI • Deepfake • Cybersecurity", fill=(0, 255, 200))

    # Face outline style
    draw.ellipse((850, 180, 1120, 450), outline=(0, 255, 200), width=6)
    draw.line((985, 230, 985, 390), fill=(255,255,255), width=4)
    draw.arc((910, 300, 1060, 420), start=0, end=180, fill=(255,255,255), width=4)

    path = f"/mnt/data/cyber_image_{i}.png"
    img.save(path)
    image_paths.append(path)

# Create slides
for idx, (title, text) in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(5, 10, 20)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.8))
    title_tf = title_box.text_frame
    run = title_tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 255, 200)

    # Body
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.2), Inches(3))
    body_tf = body_box.text_frame
    p = body_tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add image
    slide.shapes.add_picture(image_paths[idx], Inches(5), Inches(1.3), width=Inches(4.5))

# Save presentation
ppt_path = "/mnt/data/Professional_Deepfake_Cybercrime_Presentation.pptx"
prs.save(ppt_path)

print("Created:", ppt_path)
